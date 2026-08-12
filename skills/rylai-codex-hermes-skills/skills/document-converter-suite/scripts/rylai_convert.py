#!/usr/bin/env python3
"""Convert local documents with explicit, dependency-aware routes.

Original utility maintained by Rylai for Codex and Hermes skill workflows.
"""

from __future__ import annotations

import argparse
import csv
from html.parser import HTMLParser
import html
import importlib.util
import json
import os
from pathlib import Path
import re
import shutil
import subprocess
import sys
import tempfile


TEXT_INPUTS = {".txt", ".md", ".markdown"}
OFFICE_INPUTS = {".doc", ".docx", ".ppt", ".pptx", ".xls", ".xlsx", ".odt", ".odp", ".ods"}
DATA_INPUTS = {".csv", ".tsv", ".json"}


class SimpleHtmlReader(HTMLParser):
    """Extract readable blocks without executing or fetching HTML content."""

    BLOCK_TAGS = {
        "article",
        "blockquote",
        "br",
        "div",
        "footer",
        "h1",
        "h2",
        "h3",
        "h4",
        "h5",
        "h6",
        "header",
        "li",
        "main",
        "p",
        "section",
        "tr",
    }

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.parts: list[str] = []
        self.suppressed_depth = 0

    def handle_starttag(self, tag: str, attrs) -> None:
        tag = tag.lower()
        if tag in {"script", "style", "template"}:
            self.suppressed_depth += 1
        elif not self.suppressed_depth and tag in self.BLOCK_TAGS:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        tag = tag.lower()
        if tag in {"script", "style", "template"} and self.suppressed_depth:
            self.suppressed_depth -= 1
        elif not self.suppressed_depth and tag in self.BLOCK_TAGS:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self.suppressed_depth:
            self.parts.append(data)

    def plain_text(self) -> str:
        lines = []
        for raw_line in "".join(self.parts).splitlines():
            line = " ".join(raw_line.split())
            if line:
                lines.append(line)
        return "\n".join(lines).strip() + "\n"


def command_path(name: str) -> str | None:
    return shutil.which(name)


def module_available(name: str) -> bool:
    return importlib.util.find_spec(name) is not None


def markitdown_command() -> list[str] | None:
    executable = command_path("markitdown")
    if executable:
        return [executable]
    uvx = command_path("uvx")
    if uvx:
        return [uvx, "markitdown"]
    return None


def libreoffice_command() -> str | None:
    return command_path("soffice") or command_path("libreoffice")


def run_checked(command: list[str]) -> None:
    completed = subprocess.run(
        command,
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
    )
    if completed.returncode:
        message = completed.stderr.strip() or completed.stdout.strip()
        raise RuntimeError(message or f"command failed with exit code {completed.returncode}")


def ensure_source(source: Path) -> Path:
    resolved = source.expanduser().resolve()
    if not resolved.is_file():
        raise FileNotFoundError(f"source file not found: {resolved}")
    return resolved


def ensure_destination(source: Path, destination: Path, overwrite: bool) -> Path:
    resolved = destination.expanduser().resolve()
    if source == resolved:
        raise RuntimeError("source and destination must be different files")
    if resolved.exists() and not overwrite:
        raise FileExistsError(f"destination already exists: {resolved}")
    resolved.parent.mkdir(parents=True, exist_ok=True)
    return resolved


def replace_from_temp(temp_path: Path, destination: Path) -> None:
    os.replace(temp_path, destination)


def write_text_atomic(destination: Path, content: str) -> None:
    with tempfile.NamedTemporaryFile(
        mode="w",
        encoding="utf-8",
        newline="\n",
        delete=False,
        dir=destination.parent,
        prefix=f".{destination.name}.",
        suffix=".tmp",
    ) as handle:
        handle.write(content)
        temp_path = Path(handle.name)
    try:
        replace_from_temp(temp_path, destination)
    finally:
        temp_path.unlink(missing_ok=True)


def copy_atomic(source: Path, destination: Path) -> None:
    with tempfile.NamedTemporaryFile(
        delete=False,
        dir=destination.parent,
        prefix=f".{destination.name}.",
        suffix=".tmp",
    ) as handle:
        temp_path = Path(handle.name)
    try:
        shutil.copyfile(source, temp_path)
        replace_from_temp(temp_path, destination)
    finally:
        temp_path.unlink(missing_ok=True)


def read_text(source: Path) -> str:
    return source.read_text(encoding="utf-8-sig")


def cell_text(value: object) -> str:
    if value is None:
        return ""
    return str(value).replace("|", "\\|").replace("\r\n", "<br>").replace("\n", "<br>")


def markdown_table(rows: list[list[object]]) -> str:
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    normalized = [list(row) + [""] * (width - len(row)) for row in rows]
    header = normalized[0]
    body = normalized[1:]
    output = [
        "| " + " | ".join(cell_text(value) for value in header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    output.extend(
        "| " + " | ".join(cell_text(value) for value in row) + " |"
        for row in body
    )
    return "\n".join(output) + "\n"


def json_to_rows(value: object) -> list[list[object]]:
    if isinstance(value, list) and all(isinstance(item, dict) for item in value):
        keys: list[str] = []
        for item in value:
            for key in item:
                if str(key) not in keys:
                    keys.append(str(key))
        return [keys] + [[item.get(key, "") for key in keys] for item in value]
    if isinstance(value, dict):
        rows: list[list[object]] = [["key", "value"]]
        for key, item in value.items():
            rendered = json.dumps(item, ensure_ascii=False) if isinstance(item, (dict, list)) else item
            rows.append([key, rendered])
        return rows
    if isinstance(value, list):
        return [["value"]] + [[item] for item in value]
    return [["value"], [value]]


def delimited_rows(source: Path) -> list[list[str]]:
    delimiter = "\t" if source.suffix.lower() == ".tsv" else ","
    with source.open("r", encoding="utf-8-sig", newline="") as handle:
        return [list(row) for row in csv.reader(handle, delimiter=delimiter)]


def fallback_markdown(source: Path) -> str:
    extension = source.suffix.lower()
    if extension in TEXT_INPUTS:
        return read_text(source)
    if extension in {".csv", ".tsv"}:
        return markdown_table(delimited_rows(source))
    if extension == ".json":
        return markdown_table(json_to_rows(json.loads(read_text(source))))
    if extension in {".html", ".htm"}:
        parser = SimpleHtmlReader()
        parser.feed(read_text(source))
        return parser.plain_text()
    if extension == ".pdf":
        if not module_available("pypdf"):
            raise RuntimeError("PDF fallback requires the optional pypdf package")
        from pypdf import PdfReader

        reader = PdfReader(str(source))
        pages = []
        for number, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            pages.append(f"## Page {number}\n\n{text}\n")
        return "\n".join(pages)
    if extension == ".docx":
        if not module_available("docx"):
            raise RuntimeError("DOCX fallback requires the optional python-docx package")
        from docx import Documen

        document = Document(str(source))
        sections: list[str] = []
        sections.extend(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())
        for table in document.tables:
            sections.append(markdown_table([[cell.text for cell in row.cells] for row in table.rows]))
        return "\n\n".join(sections).strip() + "\n"
    if extension == ".pptx":
        if not module_available("pptx"):
            raise RuntimeError("PPTX fallback requires the optional python-pptx package")
        from pptx import Presentation

        presentation = Presentation(str(source))
        slides = []
        for number, slide in enumerate(presentation.slides, start=1):
            text_items = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_items.append(shape.text.strip())
            slides.append(f"## Slide {number}\n\n" + "\n\n".join(text_items))
        return "\n\n".join(slides).strip() + "\n"
    if extension == ".xlsx":
        if not module_available("openpyxl"):
            raise RuntimeError("XLSX fallback requires the optional openpyxl package")
        from openpyxl import load_workbook

        workbook = load_workbook(source, read_only=True, data_only=False)
        sheets = []
        for worksheet in workbook.worksheets:
            rows = [list(row) for row in worksheet.iter_rows(values_only=True)]
            while rows and not any(value not in (None, "") for value in rows[-1]):
                rows.pop()
            sheets.append(f"## Sheet: {worksheet.title}\n\n{markdown_table(rows)}")
        workbook.close()
        return "\n\n".join(sheets).strip() + "\n"
    raise RuntimeError(f"no local Markdown fallback for {extension or 'extensionless input'}")


def via_markitdown(source: Path, destination: Path) -> None:
    command = markitdown_command()
    if not command:
        raise RuntimeError("MarkItDown is unavailable; install it or install uv for uvx")
    with tempfile.NamedTemporaryFile(
        delete=False,
        dir=destination.parent,
        prefix=f".{destination.name}.",
        suffix=".md",
    ) as handle:
        temp_path = Path(handle.name)
    try:
        run_checked([*command, str(source), "-o", str(temp_path)])
        if not temp_path.is_file():
            raise RuntimeError("MarkItDown did not create an output file")
        replace_from_temp(temp_path, destination)
    finally:
        temp_path.unlink(missing_ok=True)


def via_pandoc(source: Path, destination: Path) -> None:
    pandoc = command_path("pandoc")
    if not pandoc:
        raise RuntimeError("this conversion route requires Pandoc")
    with tempfile.NamedTemporaryFile(
        delete=False,
        dir=destination.parent,
        prefix=f".{destination.name}.",
        suffix=destination.suffix,
    ) as handle:
        temp_path = Path(handle.name)
    try:
        run_checked([pandoc, str(source), "-o", str(temp_path)])
        if not temp_path.is_file() or temp_path.stat().st_size == 0:
            raise RuntimeError("Pandoc did not create a usable output file")
        replace_from_temp(temp_path, destination)
    finally:
        temp_path.unlink(missing_ok=True)


def via_libreoffice(source: Path, destination: Path) -> None:
    office = libreoffice_command()
    if not office:
        raise RuntimeError("this conversion route requires LibreOffice")
    with tempfile.TemporaryDirectory(prefix="rylai-office-") as temporary:
        temporary_dir = Path(temporary)
        run_checked(
            [
                office,
                "--headless",
                "--convert-to",
                destination.suffix.lstrip("."),
                "--outdir",
                str(temporary_dir),
                str(source),
            ]
        )
        candidates = [
            path
            for path in temporary_dir.iterdir()
            if path.is_file()
            and path.stem.casefold() == source.stem.casefold()
            and path.suffix.casefold() == destination.suffix.casefold()
        ]
        if len(candidates) != 1:
            raise RuntimeError("LibreOffice did not create the expected output")
        copy_atomic(candidates[0], destination)


def markdown_to_plain(markdown: str) -> str:
    text = re.sub(r"!\[([^\]]*)\]\([^)]+\)", r"\1", markdown)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    text = re.sub(r"(?m)^\s{0,3}#{1,6}\s*", "", text)
    text = re.sub(r"(?m)^\s*[-*+]\s+", "", text)
    text = re.sub(r"(?m)^\s*\d+[.)]\s+", "", text)
    text = re.sub(r"[*_`~]", "", text)
    return text.strip() + "\n"


def to_markdown(source: Path, destination: Path) -> str:
    if source.suffix.lower() in TEXT_INPUTS | DATA_INPUTS | {".html", ".htm"}:
        write_text_atomic(destination, fallback_markdown(source))
        return "python-standard-library"

    markitdown_error = None
    if markitdown_command():
        try:
            via_markitdown(source, destination)
            return "markitdown"
        except RuntimeError as error:
            markitdown_error = str(error)

    try:
        write_text_atomic(destination, fallback_markdown(source))
        return "python-format-fallback"
    except RuntimeError as fallback_error:
        if markitdown_error:
            raise RuntimeError(
                f"MarkItDown failed ({markitdown_error}); fallback failed ({fallback_error})"
            ) from fallback_error
        raise


def to_text(source: Path, destination: Path) -> str:
    extension = source.suffix.lower()
    if extension in {".txt"}:
        write_text_atomic(destination, read_text(source))
        return "python-text"
    if extension in {".html", ".htm"}:
        parser = SimpleHtmlReader()
        parser.feed(read_text(source))
        write_text_atomic(destination, parser.plain_text())
        return "python-html-reader"
    with tempfile.TemporaryDirectory(prefix="rylai-text-") as temporary:
        intermediate = Path(temporary) / "content.md"
        route = to_markdown(source, intermediate)
        write_text_atomic(destination, markdown_to_plain(read_text(intermediate)))
        return f"{route}-then-plain-text"


def to_csv(source: Path, destination: Path) -> str:
    extension = source.suffix.lower()
    if extension in {".csv", ".tsv"}:
        rows = delimited_rows(source)
    elif extension == ".json":
        rows = json_to_rows(json.loads(read_text(source)))
    else:
        raise RuntimeError("direct CSV output accepts CSV, TSV, or JSON; use rylai_tables.py for document tables")
    with tempfile.NamedTemporaryFile(
        mode="w",
        encoding="utf-8-sig",
        newline="",
        delete=False,
        dir=destination.parent,
        prefix=f".{destination.name}.",
        suffix=".tmp",
    ) as handle:
        csv.writer(handle).writerows(rows)
        temp_path = Path(handle.name)
    try:
        replace_from_temp(temp_path, destination)
    finally:
        temp_path.unlink(missing_ok=True)
    return "python-csv"


def to_json(source: Path, destination: Path) -> str:
    extension = source.suffix.lower()
    if extension == ".json":
        value = json.loads(read_text(source))
    elif extension in {".csv", ".tsv"}:
        delimiter = "\t" if extension == ".tsv" else ","
        with source.open("r", encoding="utf-8-sig", newline="") as handle:
            value = list(csv.DictReader(handle, delimiter=delimiter))
    else:
        raise RuntimeError("direct JSON output accepts JSON, CSV, or TSV")
    write_text_atomic(destination, json.dumps(value, ensure_ascii=False, indent=2) + "\n")
    return "python-json"


def to_html(source: Path, destination: Path) -> str:
    if source.suffix.lower() in {".html", ".htm"}:
        copy_atomic(source, destination)
        return "copy"
    if source.suffix.lower() in {".md", ".markdown"} and command_path("pandoc"):
        via_pandoc(source, destination)
        return "pandoc"
    if source.suffix.lower() in DATA_INPUTS:
        body = fallback_markdown(source)
    else:
        body = read_text(source)
    document = (
        "<!doctype html>\n"
        "<html lang=\"en\"><head><meta charset=\"utf-8\"><title>Converted document</title></head>\n"
        f"<body><pre>{html.escape(body)}</pre></body></html>\n"
    )
    write_text_atomic(destination, document)
    return "python-html-wrapper"


def convert_file(source_arg: Path, destination_arg: Path, overwrite: bool = False) -> str:
    source = ensure_source(source_arg)
    destination = ensure_destination(source, destination_arg, overwrite)
    source_extension = source.suffix.lower()
    target_extension = destination.suffix.lower()
    if not target_extension:
        raise RuntimeError("destination must have a file extension")

    if source_extension == target_extension:
        copy_atomic(source, destination)
        return "copy"
    if target_extension in {".md", ".markdown"}:
        return to_markdown(source, destination)
    if target_extension == ".txt":
        return to_text(source, destination)
    if target_extension == ".csv":
        return to_csv(source, destination)
    if target_extension == ".json":
        return to_json(source, destination)
    if target_extension in {".html", ".htm"}:
        return to_html(source, destination)
    if (
        source_extension in TEXT_INPUTS | {".html", ".htm"}
        and target_extension in {".docx", ".pdf", ".odt", ".rtf"}
    ):
        via_pandoc(source, destination)
        return "pandoc"
    if source_extension in OFFICE_INPUTS and (
        target_extension in OFFICE_INPUTS | {".pdf", ".rtf", ".html"}
    ):
        via_libreoffice(source, destination)
        return "libreoffice"
    raise RuntimeError(f"unsupported conversion route: {source_extension} -> {target_extension}")


def doctor() -> int:
    checks = [
        ("markitdown", bool(markitdown_command()), "rich semantic extraction"),
        ("pandoc", bool(command_path("pandoc")), "Markdown/HTML to document formats"),
        ("libreoffice", bool(libreoffice_command()), "Office and PDF export routes"),
        ("pypdf", module_available("pypdf"), "PDF fallback and page utilities"),
        ("python-docx", module_available("docx"), "DOCX fallback"),
        ("python-pptx", module_available("pptx"), "PPTX fallback"),
        ("openpyxl", module_available("openpyxl"), "XLSX fallback and table extraction"),
        ("pdfplumber", module_available("pdfplumber"), "PDF table extraction"),
    ]
    for name, available, purpose in checks:
        print(f"{'OK' if available else 'MISSING':7} {name:14} {purpose}")
    return 0


def batch_convert(args: argparse.Namespace) -> int:
    source_dir = args.source_dir.expanduser().resolve()
    output_dir = args.output_dir.expanduser().resolve()
    if not source_dir.is_dir():
        raise FileNotFoundError(f"source directory not found: {source_dir}")
    target_extension = "." + args.to.lower().lstrip(".")
    includes = {"." + item.lower().lstrip(".") for item in args.include}
    iterator = source_dir.rglob("*") if args.recursive else source_dir.glob("*")
    inputs = []
    for path in iterator:
        if not path.is_file():
            continue
        resolved = path.resolve()
        if output_dir == resolved.parent or output_dir in resolved.parents:
            continue
        if includes and path.suffix.lower() not in includes:
            continue
        inputs.append(path)

    failures = 0
    for source in sorted(inputs):
        relative = source.relative_to(source_dir).with_suffix(target_extension)
        destination = output_dir / relative
        try:
            route = convert_file(source, destination, args.overwrite)
            print(f"OK   {source} -> {destination} [{route}]")
        except Exception as error:
            failures += 1
            print(f"FAIL {source}: {error}", file=sys.stderr)
    print(f"Processed: {len(inputs)}; Failed: {failures}")
    return 1 if failures else 0


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description=__doc__)
    subparsers = parser.add_subparsers(dest="command", required=True)

    subparsers.add_parser("doctor", help="Report available conversion dependencies")

    convert = subparsers.add_parser("convert", help="Convert one local file")
    convert.add_argument("source", type=Path)
    convert.add_argument("destination", type=Path)
    convert.add_argument("--overwrite", action="store_true")

    batch = subparsers.add_parser("batch", help="Convert files under a directory")
    batch.add_argument("source_dir", type=Path)
    batch.add_argument("output_dir", type=Path)
    batch.add_argument("--to", required=True, help="Target extension")
    batch.add_argument("--recursive", action="store_true")
    batch.add_argument("--include", action="append", default=[], help="Source extension; repeatable")
    batch.add_argument("--overwrite", action="store_true")
    return parser


def main() -> int:
    parser = build_parser()
    args = parser.parse_args()
    try:
        if args.command == "doctor":
            return doctor()
        if args.command == "convert":
            route = convert_file(args.source, args.destination, args.overwrite)
            print(f"Wrote {args.destination} using {route}")
            return 0
        if args.command == "batch":
            return batch_convert(args)
    except Exception as error:
        print(f"ERROR: {error}", file=sys.stderr)
        return 1
    parser.error("unknown command")
    return 2


if __name__ == "__main__":
    raise SystemExit(main())
